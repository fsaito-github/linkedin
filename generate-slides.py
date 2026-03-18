"""
Generates a PowerPoint presentation comparing specialized vs generic agent outputs
for LinkedIn post generation about GitHub Copilot CLI.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Colors
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_ORANGE = RGBColor(0xF9, 0x73, 0x16)
ACCENT_BLUE = RGBColor(0x1F, 0x6F, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_GRAY = RGBColor(0xC9, 0xD1, 0xD9)
CARD_BG = RGBColor(0x16, 0x1B, 0x22)
CARD_BORDER = RGBColor(0x30, 0x36, 0x3D)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=12,
                          default_color=WHITE, font_name="Segoe UI"):
    """lines = list of (text, color, bold) tuples"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color if color else default_color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(2)
    return txBox


def add_card(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, BG_DARK)

add_accent_bar(slide1, 0, 0, 0.06, 9, ACCENT_PURPLE)

add_textbox(slide1, 1, 1.5, 14, 1, "AGENTE ESPECIALIZADO vs AGENTE GENÉRICO",
            font_size=40, color=WHITE, bold=True)

add_textbox(slide1, 1, 2.8, 14, 1,
            "Como instruções especializadas transformam a qualidade do output de IA",
            font_size=22, color=ACCENT_PURPLE)

add_multiline_textbox(slide1, 1, 4.5, 14, 3, [
    ("Mesmo prompt. Mesmo tema. Mesmo pedido.", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("A diferença? 266 linhas de instruções especializadas no .agent.md", WHITE, True),
    ("que ensinam o agente a formatar para LinkedIn, usar Unicode bold,", GRAY, False),
    ("criar hooks eficazes e seguir um checklist de qualidade.", GRAY, False),
    ("", WHITE, False),
    ("O resultado é a diferença entre um post publicável e um rascunho.", ACCENT_PURPLE, True),
], font_size=18)

add_textbox(slide1, 1, 7.5, 14, 0.5,
            "Tema: GitHub Copilot CLI  ·  Tom: Provocativo  ·  ~500 palavras  ·  PT-BR",
            font_size=14, color=GRAY)


# ============================================================
# SLIDE 2: The Prompts (side by side)
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

add_accent_bar(slide2, 0, 0, 0.06, 9, ACCENT_PURPLE)
add_textbox(slide2, 0.8, 0.3, 14, 0.7, "OS PROMPTS ENVIADOS", font_size=32, color=WHITE, bold=True)
add_textbox(slide2, 0.8, 0.9, 14, 0.5,
            "Ambos receberam instruções quase idênticas — mesmo tema, tom, tamanho e audiência",
            font_size=16, color=GRAY)

# Left card - Specialized
add_card(slide2, 0.5, 1.6, 7.2, 6.8)
add_accent_bar(slide2, 0.5, 1.6, 7.2, 0.06, ACCENT_PURPLE)
add_textbox(slide2, 0.8, 1.8, 6, 0.5, "🤖  @linkedin-tech-post  (Claude Sonnet)",
            font_size=16, color=ACCENT_PURPLE, bold=True)

prompt_spec = [
    ("Crie um post sobre GitHub Copilot CLI para LinkedIn.", WHITE, False),
    ("", WHITE, False),
    ("tema: GitHub Copilot CLI", LIGHT_GRAY, False),
    ("contexto: Uso diário no terminal para debug,", LIGHT_GRAY, False),
    ("  gerar scripts e explorar codebases.", LIGHT_GRAY, False),
    ("tom: provocativo", LIGHT_GRAY, False),
    ("tamanho: Longo (~2500 chars / ~500 palavras)", LIGHT_GRAY, False),
    ("idiomas: pt (apenas PT-BR)", LIGHT_GRAY, False),
    ("audiencia: devs e tech leads", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("Ângulo: \"Copilot saiu do editor e foi pro", WHITE, False),
    ("terminal — e isso muda tudo.\"", WHITE, False),
    ("", WHITE, False),
    ("+ 266 linhas de system prompt (.agent.md)", ACCENT_PURPLE, True),
    ("  → Formatação Unicode bold para LinkedIn", ACCENT_PURPLE, False),
    ("  → Estrutura: Hook → Corpo → CTA → Hashtags", ACCENT_PURPLE, False),
    ("  → Anti-padrões e checklist de qualidade", ACCENT_PURPLE, False),
    ("  → Exemplos de hooks eficazes", ACCENT_PURPLE, False),
    ("  → Tom e voz definidos", ACCENT_PURPLE, False),
]
add_multiline_textbox(slide2, 0.8, 2.4, 6.6, 5.5, prompt_spec, font_size=12, font_name="Consolas")

# Right card - Generic
add_card(slide2, 8.3, 1.6, 7.2, 6.8)
add_accent_bar(slide2, 8.3, 1.6, 7.2, 0.06, GRAY)
add_textbox(slide2, 8.6, 1.8, 6, 0.5, "🔧  Agente Genérico  (GPT 5.1 mini)",
            font_size=16, color=GRAY, bold=True)

prompt_gen = [
    ("Crie um post de LinkedIn sobre GitHub Copilot", WHITE, False),
    ("CLI com ~500 palavras (~2500 chars) em PT-BR.", WHITE, False),
    ("", WHITE, False),
    ("O post deve ser sobre como o GitHub Copilot", LIGHT_GRAY, False),
    ("CLI transforma o workflow de desenvolvedores.", LIGHT_GRAY, False),
    ("O autor usa diariamente para debug, gerar", LIGHT_GRAY, False),
    ("scripts e explorar codebases.", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("Tom: provocativo.", LIGHT_GRAY, False),
    ("Audiência: devs e tech leads.", LIGHT_GRAY, False),
    ("Idioma: Português BR.", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("IMPORTANTE: Gere APENAS o texto do post", WHITE, False),
    ("pronto para colar no LinkedIn.", WHITE, False),
    ("", WHITE, False),
    ("Sem system prompt especializado", ACCENT_RED, True),
    ("  → Sem regras de formatação LinkedIn", ACCENT_RED, False),
    ("  → Sem estrutura obrigatória", ACCENT_RED, False),
    ("  → Sem anti-padrões definidos", ACCENT_RED, False),
    ("  → Sem exemplos de referência", ACCENT_RED, False),
]
add_multiline_textbox(slide2, 8.6, 2.4, 6.6, 5.5, prompt_gen, font_size=12, font_name="Consolas")


# ============================================================
# SLIDE 3: Results - Specialized
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)

add_accent_bar(slide3, 0, 0, 0.06, 9, ACCENT_PURPLE)
add_textbox(slide3, 0.8, 0.3, 10, 0.7, "RESULTADO: AGENTE ESPECIALIZADO", font_size=32, color=WHITE, bold=True)
add_textbox(slide3, 0.8, 0.9, 5, 0.4, "@linkedin-tech-post  ·  Claude Sonnet  ·  24s  ·  3608 chars",
            font_size=14, color=ACCENT_PURPLE)

# Post preview (left)
add_card(slide3, 0.5, 1.5, 9.5, 7)

post_spec_lines = [
    ("Copilot no editor é útil.", WHITE, True),
    ("Copilot no terminal é onde ele começa a mudar", WHITE, True),
    ("o jeito que dev trabalha.", WHITE, True),
    ("", WHITE, False),
    ("A maioria das discussões sobre GitHub Copilot", LIGHT_GRAY, False),
    ("ainda gira em torno de autocomplete, snippets", LIGHT_GRAY, False),
    ("e sugestão de código no IDE.", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("O salto real acontece quando o Copilot sai do", LIGHT_GRAY, False),
    ("editor e entra no terminal.", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("→ erro quebrando build", ACCENT_PURPLE, False),
    ("→ comando que você não lembra", ACCENT_PURPLE, False),
    ("→ script que precisa funcionar agora", ACCENT_PURPLE, False),
    ("→ codebase nova pra entender rápido", ACCENT_PURPLE, False),
    ("", WHITE, False),
    ("𝗢 𝗺𝗲𝘂 𝘄𝗼𝗿𝗸𝗳𝗹𝗼𝘄 𝗮𝗻𝘁𝗶𝗴𝗼 𝗲𝗿𝗮 𝗲𝘀𝘀𝗲", WHITE, True),
    ("→ ler erro → Google → 5 links → adaptar → testar", GRAY, False),
    ("", WHITE, False),
    ("𝗗𝗲𝗯𝘂𝗴 · 𝗦𝗰𝗿𝗶𝗽𝘁𝘀 · 𝗘𝘅𝗽𝗹𝗼𝗿𝗮𝗿 𝗰𝗼𝗱𝗲𝗯𝗮𝘀𝗲𝘀", ACCENT_PURPLE, True),
    ("   [3 seções detalhadas com exemplos práticos]", GRAY, False),
    ("", WHITE, False),
    ("Copilot no editor ajuda a produzir código.", LIGHT_GRAY, False),
    ("Copilot CLI ajuda a destravar trabalho.", WHITE, True),
    ("", WHITE, False),
    ("Você já testou Copilot CLI no seu fluxo real?", ACCENT_GREEN, True),
    ("", WHITE, False),
    ("#GitHubCopilot #CopilotCLI #DevTools #AI", ACCENT_PURPLE, False),
    ("#DeveloperExperience", ACCENT_PURPLE, False),
]
add_multiline_textbox(slide3, 0.8, 1.7, 8.8, 6.5, post_spec_lines, font_size=12)

# Annotations (right)
add_card(slide3, 10.3, 1.5, 5.2, 7)
annotations = [
    ("✅  O QUE O AGENTE ACERTOU", ACCENT_GREEN, True),
    ("", WHITE, False),
    ("Hook forte", WHITE, True),
    ("Primeira linha para o scroll, sem clichê", GRAY, False),
    ("", WHITE, False),
    ("Unicode bold nos subtítulos", WHITE, True),
    ("𝗦𝗲𝗰̧𝗼̃𝗲𝘀 formatadas para LinkedIn", GRAY, False),
    ("", WHITE, False),
    ("→ como bullet points", WHITE, True),
    ("Não usa - ou • (LinkedIn não renderiza bem)", GRAY, False),
    ("", WHITE, False),
    ("Parágrafos de 1-3 linhas", WHITE, True),
    ("Escaneável no feed mobile", GRAY, False),
    ("", WHITE, False),
    ("Estrutura completa", WHITE, True),
    ("Hook → Antes/Depois → 3 seções → CTA → Hashtags", GRAY, False),
    ("", WHITE, False),
    ("CTA natural", WHITE, True),
    ("Pergunta aberta que convida experiência", GRAY, False),
    ("", WHITE, False),
    ("5 hashtags relevantes", WHITE, True),
    ("Mix de amplas (#AI) e específicas (#CopilotCLI)", GRAY, False),
    ("", WHITE, False),
    ("Tom provocativo sem ser arrogante", WHITE, True),
    ("\"A busca tradicional virou hábito, não vantagem\"", GRAY, False),
]
add_multiline_textbox(slide3, 10.6, 1.7, 4.6, 6.5, annotations, font_size=11)


# ============================================================
# SLIDE 4: Results - Generic
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)

add_accent_bar(slide4, 0, 0, 0.06, 9, ACCENT_RED)
add_textbox(slide4, 0.8, 0.3, 10, 0.7, "RESULTADO: AGENTE GENÉRICO", font_size=32, color=WHITE, bold=True)
add_textbox(slide4, 0.8, 0.9, 5, 0.4, "GPT 5.1 mini  ·  6s  ·  1495 chars (pedido: 2500)",
            font_size=14, color=ACCENT_RED)

# Post preview (left)
add_card(slide4, 0.5, 1.5, 9.5, 7)

post_gen_lines = [
    ("Debugar no terminal já era meu pesadelo. Agora,", WHITE, False),
    ("o GitHub Copilot CLI virou meu copiloto literal:", WHITE, False),
    ("aperto uma tecla e ele me mostra o que está", WHITE, False),
    ("quebrando, gera scripts repetitivos e até me", WHITE, False),
    ("ajuda a entender codebases inteiras em segundos.", WHITE, False),
    ("Antes eu precisava abrir várias abas, esperar", WHITE, False),
    ("IDEs carregarem e ficar pulando entre contextos", WHITE, False),
    ("— hoje, tudo acontece na linha de comando...", WHITE, False),
    ("", WHITE, False),
    ("O Copilot CLI me força a pensar como um prompt", WHITE, False),
    ("engineer: descrevo o que quero resolver e ele", WHITE, False),
    ("sugere o próximo passo. O resultado? Menos tempo", WHITE, False),
    ("perdido em configurações e mais foco em entregar", WHITE, False),
    ("valor. Preciso de um script de deploy? Ele gera.", WHITE, False),
    ("Quero entender como a função X é usada em todo", WHITE, False),
    ("o projeto? Ele mapeia...", WHITE, False),
    ("", WHITE, False),
    ("E não é mágica isolada — é um workflow. Integro", WHITE, False),
    ("Copilot CLI ao meu dia a dia: uso para descobrir", WHITE, False),
    ("dependências, gerar testes, verificar logs...", WHITE, False),
    ("", WHITE, False),
    ("Se você ainda não usa, experimente: sente no", WHITE, False),
    ("terminal, descreva seu problema e veja o Copilot", WHITE, False),
    ("responder em minutos. O futuro do desenvolvimento", WHITE, False),
    ("está no prompt, e quem dominar essa conversa vai", WHITE, False),
    ("liderar equipes mais rápidas e inteligentes.", WHITE, False),
]
add_multiline_textbox(slide4, 0.8, 1.7, 8.8, 6.5, post_gen_lines, font_size=12)

# Annotations (right)
add_card(slide4, 10.3, 1.5, 5.2, 7)
problems = [
    ("❌  O QUE O AGENTE ERROU", ACCENT_RED, True),
    ("", WHITE, False),
    ("Tamanho insuficiente", WHITE, True),
    ("1495 chars vs 2500 pedidos (40% menor)", GRAY, False),
    ("", WHITE, False),
    ("Sem formatação LinkedIn", WHITE, True),
    ("Sem Unicode bold, sem → , sem quebras visuais", GRAY, False),
    ("", WHITE, False),
    ("4 parágrafos densos", WHITE, True),
    ("Textos longos = \"muro de texto\" no feed mobile", GRAY, False),
    ("", WHITE, False),
    ("Sem hashtags", WHITE, True),
    ("Post invisível na busca do LinkedIn", GRAY, False),
    ("", WHITE, False),
    ("Sem CTA", WHITE, True),
    ("Não convida interação — post morre sem engajamento", GRAY, False),
    ("", WHITE, False),
    ("Sem estrutura definida", WHITE, True),
    ("Não segue Hook → Corpo → Fechamento", GRAY, False),
    ("", WHITE, False),
    ("Tom de artigo de blog", WHITE, True),
    ("\"Quem dominar essa conversa vai liderar...\"", GRAY, False),
    ("Soa como propaganda, não experiência pessoal", GRAY, False),
    ("", WHITE, False),
    ("Hook fraco", WHITE, True),
    ("\"Debugar no terminal já era meu pesadelo\"", GRAY, False),
    ("Genérico — poderia ser de qualquer ferramenta", GRAY, False),
]
add_multiline_textbox(slide4, 10.6, 1.7, 4.6, 6.5, problems, font_size=11)


# ============================================================
# SLIDE 5: Side-by-side comparison table
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)

add_accent_bar(slide5, 0, 0, 0.06, 9, ACCENT_PURPLE)
add_textbox(slide5, 0.8, 0.3, 14, 0.7, "COMPARATIVO DIRETO", font_size=32, color=WHITE, bold=True)
add_textbox(slide5, 0.8, 0.9, 14, 0.5,
            "Mesmo pedido, resultados completamente diferentes",
            font_size=16, color=GRAY)

# Table rows
criteria = [
    ("Critério",            "Especializado",                    "Genérico"),
    ("Tempo",               "24 segundos",                      "6 segundos"),
    ("Tamanho",             "3608 chars ✅",                    "1495 chars ❌ (40% do pedido)"),
    ("Unicode bold",        "✅ Subtítulos formatados",         "❌ Não usa"),
    ("Bullet points →",     "✅ Consistente",                   "❌ Texto corrido"),
    ("Parágrafos curtos",   "✅ 1-3 linhas",                    "❌ 4-8 linhas (muro de texto)"),
    ("Hook",                "✅ Forte e específico",             "⚠️ Genérico"),
    ("Estrutura",           "✅ Hook→Corpo→CTA→Hashtags",       "❌ Sem estrutura definida"),
    ("Hashtags",            "✅ 5 relevantes",                  "❌ Nenhuma"),
    ("CTA",                 "✅ Pergunta aberta natural",       "❌ Tom imperativo"),
    ("Escaneabilidade",     "✅ Alta (feed mobile)",             "❌ Baixa (muro de texto)"),
    ("Pronto p/ publicar?", "✅ SIM — copiar e colar",          "❌ NÃO — precisa reescrever"),
]

row_h = 0.5
start_y = 1.6
col_widths = [2.8, 5.5, 6.5]
col_starts = [0.6, 3.5, 9.1]

for i, (col0, col1, col2) in enumerate(criteria):
    y = start_y + i * row_h
    is_header = (i == 0)
    bg = CARD_BG if i % 2 == 1 else RGBColor(0x1C, 0x21, 0x28)

    if is_header:
        bg = RGBColor(0x21, 0x26, 0x2D)

    for j, (text, cw, cs) in enumerate([(col0, col_widths[0], col_starts[0]),
                                         (col1, col_widths[1], col_starts[1]),
                                         (col2, col_widths[2], col_starts[2])]):
        card = add_card(slide5, cs, y, cw, row_h, bg)
        color = WHITE if is_header else (LIGHT_GRAY if j == 0 else WHITE)
        if not is_header and j == 1 and "✅" in text:
            color = ACCENT_GREEN
        elif not is_header and j == 2 and "❌" in text:
            color = ACCENT_RED
        elif not is_header and j == 2 and "⚠️" in text:
            color = ACCENT_ORANGE

        add_textbox(slide5, cs + 0.15, y + 0.08, cw - 0.3, row_h - 0.15, text,
                    font_size=13 if not is_header else 14,
                    color=color, bold=is_header)


# ============================================================
# SLIDE 6: Key Insight / Conclusion
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_DARK)

add_accent_bar(slide6, 0, 0, 0.06, 9, ACCENT_PURPLE)

add_textbox(slide6, 1, 1, 14, 1, "POR QUE ISSO ACONTECE?", font_size=40, color=WHITE, bold=True)

add_multiline_textbox(slide6, 1, 2.3, 14, 1.5, [
    ("O prompt do usuário foi quase idêntico nos dois casos.", LIGHT_GRAY, False),
    ("A diferença está no system prompt — as instruções do agente.", WHITE, True),
], font_size=22)

# Three cards
cards_data = [
    ("AGENTE ESPECIALIZADO", ACCENT_PURPLE, [
        "266 linhas de instruções em .agent.md",
        "Regras de formatação LinkedIn (Unicode bold, →)",
        "Estrutura obrigatória (Hook → CTA → Hashtags)",
        "Anti-padrões explícitos (o que NUNCA fazer)",
        "Checklist de qualidade com 8 critérios",
        "Exemplos de hooks eficazes como referência",
        "Tom e voz definidos com exemplos",
    ]),
    ("AGENTE GENÉRICO", GRAY, [
        "Apenas o prompt do usuário",
        "Sem conhecimento de formatação LinkedIn",
        "Sem estrutura obrigatória",
        "Sem anti-padrões definidos",
        "Sem checklist de qualidade",
        "Sem exemplos de referência",
        "Tom genérico de \"assistente útil\"",
    ]),
    ("CONCLUSÃO", ACCENT_GREEN, [
        "Instruções especializadas = consistência",
        "O agente entrega post publicável, não rascunho",
        "Elimina retrabalho de formatação manual",
        "Garante padrão em posts diferentes",
        "É a diferença entre automação e assistência",
        "",
        "System prompt > prompt do usuário",
    ]),
]

for i, (title, color, items) in enumerate(cards_data):
    x = 0.7 + i * 5.1
    add_card(slide6, x, 4, 4.7, 4.2)
    add_accent_bar(slide6, x, 4, 4.7, 0.06, color)
    add_textbox(slide6, x + 0.3, 4.2, 4, 0.5, title, font_size=16, color=color, bold=True)

    lines = [(item, LIGHT_GRAY if item else WHITE, item.startswith("System")) for item in items]
    # Make last item bold if it's the conclusion card
    if i == 2 and lines:
        last = lines[-1]
        lines[-1] = (last[0], ACCENT_GREEN, True)

    add_multiline_textbox(slide6, x + 0.3, 4.8, 4.1, 3.2, lines, font_size=13)


# Save
output_path = os.path.join(r"C:\Users\fabiosaito\linkedin", "comparativo-agente-especializado-vs-generico.pptx")
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
